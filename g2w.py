import os
import numpy as np
import cv2
import shutil
import zipfile
from alive_progress import alive_bar
from pptx import Presentation
from pptx.util import Pt
from pptx.shapes.picture import Picture
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
import json
from PIL import Image
import io

def crop_image(im):
    img = cv2.imread(im)
    img = img[40:-70,0:-1]
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    gray = 255*(gray < 128).astype(np.uint8)
    gray = cv2.morphologyEx(gray, cv2.MORPH_OPEN, np.ones((2, 2), dtype=np.uint8))
    coords = cv2.findNonZero(gray)
    x, y, w, h = cv2.boundingRect(coords)
    w = w + x + 100
    rect = img[y:y+h, 0:w]
    cv2.imwrite(im, rect)
    
def generate_metadata():
    print("1. Extracting images and generating metadata.")
    os.mkdir("temp")
    # Find all melodies in /hymn folder
    melodies = {}
    for subdir, dirs, files in os.walk(r'hymns'):
        for filename in files:
            filepath = subdir + os.sep + filename
            if filepath.endswith(".ppt") and "Melody" in filepath:
                melodies[filename.replace('.ppt', '')] = filepath
   
    with alive_bar(len(melodies)) as bar:
        for key, value in melodies.items():
            d = "temp/" + key
            os.mkdir(d)
            os.mkdir(d + "/images")
            # Fetch title slide information from powerpoint
            prs = Presentation(value)
            shapes = prs.slides[0].shapes
            firstSlide = 0
            if len(shapes) < 2:
                shapes = prs.slides[1].shapes
                firstSlide = 1
            hymn_name = shapes[0].text
            hymn_number = shapes[1].text
            hymn_credits = shapes[2].text
            data = {}
            data['name'] = hymn_name.strip()
            data['number'] = hymn_number.strip()
            data['credits'] = hymn_credits.strip()
            # Export images from each slide
            for x in range(firstSlide + 1, len(prs.slides)):
                f = open(d + "/images/" + str(x - firstSlide) + ".png", "wb")
                f.write(prs.slides[x].shapes[0].image.blob)
                f.close()
            with open(d + "/metadata.json", 'w') as outfile:
                json.dump(data, outfile)
            bar()
    print("\n")
    
def crop_images():
    print("2. Cropping melody images.")
    images = []
    for subdir, dirs, files in os.walk(r'temp'):
        for filename in files:
            filepath = subdir + os.sep + filename
            if filepath.endswith(".png"):
                images.append(filepath)
                    
    with alive_bar(len(images)) as bar:
        for file in images:
            crop_image(file)
            bar()
    print("\n")

def create_presentations():
    print("3. Generating final hymns.")
    temp = []
    for path in os.listdir("temp"):
        temp.append(path)
        
    with alive_bar(len(temp)) as bar:
        for path in temp:
            bar.text(path + " (Title Slide)")
            d = "temp/" + path
            prs = Presentation()
            # Configure as 16:9 Aspect Ratio
            prs.slide_height = 6858000
            prs.slide_width = 12192000
            shapes = prs.slides.add_slide(prs.slide_layouts[6]).shapes
            with open(d + "/metadata.json", 'r') as f:
                title_assets =  json.loads(f.read())
            # Generate title textbox
            name_frame = shapes.add_textbox(914400, 2130426, 10363200, 1470525).text_frame
            name_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            name = name_frame.paragraphs[0]
            name.text = title_assets['name']
            name.font.size = Pt(60)
            name.font.bold = True
            name.alignment = PP_ALIGN.CENTER
            name_frame.fit_text('Calibri', 60, True)
            # Generate GTG number textbox
            number_frame = shapes.add_textbox(2341418, 6928, 7772400, 735013).text_frame
            number_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            number = number_frame.paragraphs[0]
            number.text = title_assets['number']
            number.font.size = Pt(30)
            number.font.italic = True
            number.alignment = PP_ALIGN.CENTER
            # Generate credits textbox
            credits_frame = shapes.add_textbox(1513609, 5688449, 9164782, 1169551).text_frame
            credits_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            credits = credits_frame.paragraphs[0]
            credits.text = title_assets['credits']
            credits.font.size = Pt(14)
            credits.font.italic = True
            credits.alignment = PP_ALIGN.CENTER
            credits_frame.fit_text('Calibri', 14)
            # Generate image slides
            images = {}
            bar.text(path + " (Images)")                    
            for subdir, dirs, files in os.walk(d):
                for filename in files:
                    filepath = subdir + os.sep + filename
                    if filepath.endswith(".png"):
                        images[int(filename.replace('.png', ''))] = filepath
            # Sort images by number
            images = dict(sorted(images.items()))
            for key, value in dict(sorted(images.items())).items():
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    picture = slide.shapes.add_picture(value, 0, 0, prs.slide_width)
                    calc_top_value = round((prs.slide_height - picture.height) / 2)
                    picture.top = calc_top_value         
            prs.save("out/" + path + "_wide.pptx")
            bar()
    print("\n")
    
def clean_up():
    print("4. Cleaning up.")
    temp = []
    for path in os.listdir("temp"):
        temp.append(path)
        
    with alive_bar(len(temp)) as bar:
        for path in temp:
            shutil.rmtree("temp/" + path)
            bar()
    shutil.rmtree("temp")
    print("\n")
    
if __name__ == "__main__":
    print("\nGlory2Wide by Noah Husby\n")
    if os.path.isdir("temp"):
        shutil.rmtree("temp")
    if os.path.isdir("out") is False:
        os.mkdir("out")
    if os.path.isdir("hymns") is False:
        os.mkdir("hymns")
        print("Please extract the Glory to God hymn pack into the \"hymns\" directory")
        quit()
    generate_metadata()
    crop_images()
    create_presentations()
    clean_up()
    print("Successfully converted hymns to wide format.\nThe exported hymns are in the \"out\" folder.\n")

    